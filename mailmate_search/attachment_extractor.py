"""Extract text content from email attachments."""

import io
import logging
from typing import Optional

logger = logging.getLogger(__name__)

# Maximum attachment size to process (10MB)
MAX_ATTACHMENT_SIZE = 10 * 1024 * 1024


def extract_text_from_attachment(
    attachment_data: bytes, content_type: str, filename: str
) -> Optional[str]:
    """
    Extract text content from an attachment.
    
    Args:
        attachment_data: Raw bytes of the attachment
        content_type: MIME content type (e.g., 'application/pdf')
        filename: Original filename of the attachment
        
    Returns:
        Extracted text as string, or None if extraction fails or unsupported type
    """
    # Check size limit
    if len(attachment_data) > MAX_ATTACHMENT_SIZE:
        logger.warning(
            f"Attachment {filename} exceeds size limit ({len(attachment_data)} bytes), skipping text extraction"
        )
        return None
    
    if not attachment_data:
        return None
    
    # Determine file type from content_type and filename
    content_type_lower = content_type.lower() if content_type else ""
    filename_lower = filename.lower() if filename else ""
    
    # PDF files
    if "pdf" in content_type_lower or filename_lower.endswith(".pdf"):
        return _extract_pdf_text(attachment_data)
    
    # Word documents (.docx)
    if (
        "wordprocessingml" in content_type_lower
        or "msword" in content_type_lower
        or filename_lower.endswith(".docx")
    ):
        return _extract_docx_text(attachment_data)
    
    # Excel files (.xlsx)
    if (
        "spreadsheetml" in content_type_lower
        or "excel" in content_type_lower
        or filename_lower.endswith(".xlsx")
        or filename_lower.endswith(".xlsm")
    ):
        return _extract_xlsx_text(attachment_data)
    
    # PowerPoint files (.pptx)
    if (
        "presentationml" in content_type_lower
        or "powerpoint" in content_type_lower
        or filename_lower.endswith(".pptx")
        or filename_lower.endswith(".pptm")
    ):
        return _extract_pptx_text(attachment_data)
    
    # Plain text files
    if "text/plain" in content_type_lower or filename_lower.endswith(".txt"):
        return _extract_text_file(attachment_data)
    
    # CSV files
    if "text/csv" in content_type_lower or filename_lower.endswith(".csv"):
        return _extract_csv_text(attachment_data)
    
    # Unsupported type
    logger.debug(f"Unsupported attachment type: {content_type} ({filename})")
    return None


def _extract_pdf_text(data: bytes) -> Optional[str]:
    """Extract text from PDF using pdfplumber."""
    try:
        import pdfplumber
        
        with pdfplumber.open(io.BytesIO(data)) as pdf:
            text_parts = []
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_parts.append(page_text)
            return "\n".join(text_parts) if text_parts else None
    except ImportError:
        logger.warning("pdfplumber not available, cannot extract PDF text")
        return None
    except Exception as e:
        logger.warning(f"Failed to extract text from PDF: {e}")
        return None


def _extract_docx_text(data: bytes) -> Optional[str]:
    """Extract text from Word document using python-docx."""
    try:
        from docx import Document
        
        doc = Document(io.BytesIO(data))
        text_parts = []
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)
        
        # Extract tables
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    text_parts.append(row_text)
        
        return "\n".join(text_parts) if text_parts else None
    except ImportError:
        logger.warning("python-docx not available, cannot extract Word document text")
        return None
    except Exception as e:
        logger.warning(f"Failed to extract text from Word document: {e}")
        return None


def _extract_xlsx_text(data: bytes) -> Optional[str]:
    """Extract text from Excel file using openpyxl."""
    try:
        from openpyxl import load_workbook
        
        workbook = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
        text_parts = []
        
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            sheet_text = [f"Sheet: {sheet_name}"]
            
            for row in sheet.iter_rows(values_only=True):
                row_values = [str(cell) if cell is not None else "" for cell in row]
                row_text = " | ".join(row_values).strip()
                if row_text:
                    sheet_text.append(row_text)
            
            if len(sheet_text) > 1:  # More than just the sheet name
                text_parts.append("\n".join(sheet_text))
        
        workbook.close()
        return "\n\n".join(text_parts) if text_parts else None
    except ImportError:
        logger.warning("openpyxl not available, cannot extract Excel text")
        return None
    except Exception as e:
        logger.warning(f"Failed to extract text from Excel file: {e}")
        return None


def _extract_pptx_text(data: bytes) -> Optional[str]:
    """Extract text from PowerPoint file using python-pptx."""
    try:
        from pptx import Presentation
        
        prs = Presentation(io.BytesIO(data))
        text_parts = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"Slide {slide_num}:"]
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            # Extract text from notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    notes_text = notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        slide_text.append(f"Notes: {notes_text}")
            
            if len(slide_text) > 1:  # More than just the slide number
                text_parts.append("\n".join(slide_text))
        
        return "\n\n".join(text_parts) if text_parts else None
    except ImportError:
        logger.warning("python-pptx not available, cannot extract PowerPoint text")
        return None
    except Exception as e:
        logger.warning(f"Failed to extract text from PowerPoint file: {e}")
        return None


def _extract_text_file(data: bytes) -> Optional[str]:
    """Extract text from plain text file with encoding detection."""
    # Try UTF-8 first
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        pass
    
    # Try common encodings
    for encoding in ["latin-1", "cp1252", "iso-8859-1"]:
        try:
            return data.decode(encoding, errors="replace")
        except (UnicodeDecodeError, LookupError):
            continue
    
    # Fallback: decode with errors replaced
    try:
        return data.decode("utf-8", errors="replace")
    except Exception as e:
        logger.warning(f"Failed to decode text file: {e}")
        return None


def _extract_csv_text(data: bytes) -> Optional[str]:
    """Extract text from CSV file."""
    # For CSV, we'll just decode as text
    # Could parse it more intelligently, but for search purposes, raw text is fine
    return _extract_text_file(data)

